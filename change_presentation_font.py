import os
import sys
import argparse
from pptx import Presentation

def change_font_in_shape(shape, font_name):
    """Рекурсивно меняет шрифт во всех текстовых элементах фигуры."""
    # 1. Обработка обычного текстового блока
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = font_name
                
    # 2. Обработка таблиц
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                if cell.text_frame:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = font_name
                            
    # 3. Рекурсивный обход сгруппированных фигур (групп)
    if hasattr(shape, "shapes"):
        for sub_shape in shape.shapes:
            change_font_in_shape(sub_shape, font_name)

def change_presentation_font(input_path, output_path=None, font_name="Montserrat"):
    """Загружает презентацию, меняет шрифты и сохраняет результат."""
    if not os.path.exists(input_path):
        print(f"Ошибка: Путь '{input_path}' не найден!")
        return False
        
    if os.path.isdir(input_path):
        print(f"Ошибка: Указанный путь '{input_path}' является ПАПКОЙ (директорией), а не файлом презентации!")
        print("Пожалуйста, укажите точный путь к файлу с расширением .pptx.")
        return False
        
    try:
        import io
        with open(input_path, "rb") as f:
            source_stream = io.BytesIO(f.read())
        prs = Presentation(source_stream)
        
        # Обходим каждый слайд
        for slide in prs.slides:

            # Обходим каждую фигуру на слайде
            for shape in slide.shapes:
                change_font_in_shape(shape, font_name)
                
        # Определяем путь сохранения
        if not output_path:
            output_path = input_path  # перезаписать исходный файл
            
        prs.save(output_path)
        print(f"Успех! Шрифт во всей презентации изменен на '{font_name}'.")
        print(f"Результат сохранен в: {output_path}")
        return True
    except PermissionError:
        print(f"\nОшибка доступа (Permission denied) для: {input_path}")
        print("=" * 60)
        print("Возможные причины:")
        print("1. Файл сейчас ОТКРЫТ в программе PowerPoint. PowerPoint блокирует файлы при редактировании.")
        print("   РЕШЕНИЕ: Закройте презентацию в PowerPoint и запустите скрипт снова.")
        print("2. Вы указали путь к папке вместо файла (хотя мы проверили это на старте).")
        print("3. У вашей учетной записи Windows нет прав на запись в эту папку/файл.")
        print("=" * 60)
        return False
    except Exception as e:
        print(f"Произошла ошибка при изменении шрифта: {e}")
        return False


def interactive_mode():
    """Интерактивный опрос пользователя при запуске без аргументов."""
    print("=== УТИЛИТА СМЕНЫ ШРИФТА В ПРЕЗЕНТАЦИИ PPTX ===")
    print("Этот скрипт заменит все шрифты в текстовых блоках и таблицах на выбранный вами шрифт.")
    print("-" * 50)
    
    # 1. Запрос пути к файлу
    while True:
        input_path = input("Введите путь к файлу презентации .pptx: ").strip().strip('"\'')
        if not input_path:
            print("Путь не может быть пустым.")
            continue
        if not os.path.exists(input_path):
            print("Файл по указанному пути не найден. Попробуйте еще раз.")
            continue
        if os.path.isdir(input_path):
            print("Указанный путь является ПАПКОЙ. Пожалуйста, укажите путь к конкретному файлу презентации (например, 'Приложение.pptx').")
            continue
        if not input_path.lower().endswith('.pptx'):
            print("Скрипт работает только с файлами в формате .pptx.")
            continue
        break

        
    # 2. Запрос шрифта
    font_name = input("Введите название шрифта (по умолчанию: Montserrat): ").strip()
    if not font_name:
        font_name = "Montserrat"
        
    # 3. Запрос пути для сохранения
    save_choice = input("Перезаписать исходный файл? (y/n, по умолчанию: n): ").strip().lower()
    
    if save_choice == 'y':
        output_path = input_path
    else:
        dir_name, file_name = os.path.split(input_path)
        name, ext = os.path.splitext(file_name)
        default_out = os.path.join(dir_name, f"{name}_fixed{ext}")
        
        output_path = input(f"Введите путь для сохранения (по умолчанию: {default_out}): ").strip().strip('"\'')
        if not output_path:
            output_path = default_out
            
    print("\nВыполняется обработка презентации...")
    change_presentation_font(input_path, output_path, font_name)

def main():
    # Настройка парсера аргументов командной строки
    parser = argparse.ArgumentParser(description="Утилита для смены шрифта во всех слайдах презентации PPTX.")
    parser.add_argument("-i", "--input", help="Путь к исходному файлу презентации .pptx")
    parser.add_argument("-o", "--output", help="Путь к выходному файлу (по умолчанию перезаписывает исходный)")
    parser.add_argument("-f", "--font", default="Montserrat", help="Имя шрифта (по умолчанию: Montserrat)")
    
    args = parser.parse_args()
    
    # Если аргументы командной строки не переданы, запускаем интерактивный CLI режим
    if len(sys.argv) == 1:
        interactive_mode()
    else:
        if not args.input:
            parser.print_help()
            sys.exit(1)
        change_presentation_font(args.input, args.output, args.font)

if __name__ == "__main__":
    main()
