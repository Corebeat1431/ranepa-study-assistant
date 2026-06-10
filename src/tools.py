import os
import urllib.request
import urllib.parse
import json
import xml.etree.ElementTree as ET
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import pypdf
from duckduckgo_search import DDGS
from crewai.tools import tool

# Создаем папки для выходных файлов, если их нет
os.makedirs("output", exist_ok=True)

# =====================================================================
# ИНСТРУМЕНТ 1. ПОИСК НА ARXIV
# =====================================================================
@tool
def search_arxiv(query: str) -> str:
    """Ищет научные статьи на портале arXiv.org по заданной теме.
    Возвращает заголовки, авторов и аннотации статей.
    Запрос (query) должен быть на АНГЛИЙСКОМ языке."""
    formatted_query = query.strip().replace(" ", "+")
    url = f"http://export.arxiv.org/api/query?search_query=all:{formatted_query}&max_results=3&sortBy=submittedDate&sortOrder=descending"
    
    try:
        with urllib.request.urlopen(url) as response:
            xml_data = response.read()
        
        root = ET.fromstring(xml_data)
        results = []
        ns = {'atom': 'http://www.w3.org/2005/Atom'}
        
        for entry in root.findall('atom:entry', ns):
            title = entry.find('atom:title', ns).text.strip().replace('\n', ' ')
            summary = entry.find('atom:summary', ns).text.strip().replace('\n', ' ')
            published = entry.find('atom:published', ns).text[:10]
            authors = [author.find('atom:name', ns).text for author in entry.findall('atom:author', ns)]
            
            results.append(
                f"Название: {title}\n"
                f"Авторы: {', '.join(authors)}\n"
                f"Опубликовано: {published}\n"
                f"Аннотация: {summary}\n"
                f"--------------------------------------------------"
            )
        
        if not results:
            return f"По запросу '{query}' научных статей не найдено."
        
        return "\n".join(results)
        
    except Exception as e:
        return f"Ошибка при поиске на arXiv: {e}"


# =====================================================================
# ИНСТРУМЕНТ 2. ГЕНЕРАТОР ГРАФИКОВ (MATPLOTLIB)
# =====================================================================
@tool
def generate_chart(chart_type: str, title: str, labels_str: str, values_str: str, filename: str) -> str:
    """Генерирует график (круговой, линейный или столбчатый) и сохраняет его в папку output/.
    
    Параметры:
    - chart_type: тип графика (строго 'bar' для столбчатого, 'line' для линейного, 'pie' для кругового).
    - title: заголовок графика на русском языке.
    - labels_str: список подписей через запятую (например: '2023, 2024, 2025' или 'Транспорт, Экология, ЖКХ').
    - values_str: список соответствующих числовых значений через запятую (например: '100, 150, 200' или '45.2, 30.1, 24.7').
    - filename: имя сохраняемого файла картинки (строго с расширением .png, например: 'budget_chart.png').
    """
    try:
        labels = [x.strip() for x in labels_str.split(',')]
        values = [float(x.strip()) for x in values_str.split(',')]
        
        if len(labels) != len(values):
            return "Ошибка: количество подписей и значений должно совпадать."
        
        primary_color = '#8B0000' # Темно-красный бордовый
        secondary_colors = ['#8B0000', '#A52A2A', '#D2691E', '#7F7F7F', '#A9A9A9', '#D3D3D3']
        
        plt.figure(figsize=(8, 5))
        plt.title(title, fontsize=14, fontweight='bold', pad=15, color='#333333')
        
        if chart_type.lower() == 'bar':
            bars = plt.bar(labels, values, color=primary_color, edgecolor='#555555', width=0.6)
            plt.grid(axis='y', linestyle='--', alpha=0.5)
            for bar in bars:
                height = bar.get_height()
                plt.text(bar.get_x() + bar.get_width()/2.0, height, f'{height:g}', 
                         ha='center', va='bottom', fontsize=10, fontweight='bold')
                
        elif chart_type.lower() == 'line':
            plt.plot(labels, values, marker='o', linewidth=3, color=primary_color, markersize=8)
            plt.grid(True, linestyle='--', alpha=0.5)
            for x, y in zip(labels, values):
                plt.text(x, y, f'{y:g}', ha='center', va='bottom', fontsize=10, fontweight='bold')
                
        elif chart_type.lower() == 'pie':
            plt.pie(values, labels=labels, autopct='%1.1f%%', 
                    colors=secondary_colors[:len(values)], startangle=140, 
                    textprops={'fontsize': 11, 'weight': 'bold'})
            
        else:
            return f"Ошибка: неподдерживаемый тип графика '{chart_type}'. Используйте: bar, line, pie."
        
        plt.tight_layout()
        save_path = os.path.join("output", filename)
        plt.savefig(save_path, dpi=150)
        plt.close()
        
        return f"График успешно создан и сохранен как: {save_path}"
        
    except Exception as e:
        plt.close()
        return f"Ошибка при генерации графика: {e}"


# =====================================================================
# ИНСТРУМЕНТ 3. ЧТЕНИЕ ФАЙЛОВ ПОЛЬЗОВАТЕЛЯ
# =====================================================================
@tool
def read_text_file(file_path: str) -> str:
    """Считывает содержимое текстового файла (.txt или .md) с диска. 
    Полезно для загрузки готового текста ВКР пользователя.
    Входной параметр file_path — путь к файлу."""
    if not os.path.exists(file_path):
        return f"Ошибка: файл '{file_path}' не найден."
    
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            content = f.read()
        return content
    except Exception as e:
        return f"Не удалось прочитать файл: {e}"


# =====================================================================
# ИНСТРУМЕНТ 4. ПОИСК В ОБЫЧНОМ ИНТЕРНЕТЕ (DUCKDUCKGO)
# =====================================================================
@tool
def search_web(query: str) -> str:
    """Ищет информацию в обычном интернете (через DuckDuckGo).
    Возвращает заголовки, ссылки и краткие описания веб-страниц на русском языке.
    Идеяльно подходит для поиска актуальных новостей, муниципальных данных и фактов."""
    try:
        with DDGS() as ddgs:
            results = list(ddgs.text(query, max_results=5, region="ru-ru"))
            
        formatted_results = []
        for r in results:
            formatted_results.append(
                f"Заголовок: {r['title']}\n"
                f"Ссылка: {r['href']}\n"
                f"Описание: {r['body']}\n"
                f"--------------------------------------------------"
            )
        
        if not formatted_results:
            return f"По запросу '{query}' ничего не найдено."
            
        return "\n".join(formatted_results)
    except Exception as e:
        return f"Ошибка при поиске в интернете: {e}"


# =====================================================================
# ИНСТРУМЕНТ 5. ЧТЕНИЕ PDF-ФАЙЛОВ
# =====================================================================
@tool
def read_pdf_file(file_path: str) -> str:
    """Считывает текстовое содержимое PDF-файла (например, методического пособия) 
    и возвращает весь текст для анализа агентами.
    Входной параметр file_path — путь к файлу."""
    if not os.path.exists(file_path):
        return f"Ошибка: file '{file_path}' не найден."
    
    try:
        reader = pypdf.PdfReader(file_path)
        text_content = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            if text:
                text_content.append(f"--- Страница {i+1} ---\n{text}")
        
        if not text_content:
            return f"В PDF-файле '{file_path}' не обнаружен текст (возможно, это сканированный документ без распознанного слоя)."
            
        return "\n\n".join(text_content)
    except Exception as e:
        return f"Не удалось прочитать PDF-файл: {e}"


# =====================================================================
# ИНСТРУМЕНТ 5.5. ЧТЕНИЕ WORD-ФАЙЛОВ (.DOCX)
# =====================================================================
@tool
def read_docx_file(file_path: str) -> str:
    """Считывает текстовое содержимое файла Microsoft Word (.docx).
    Извлекает текст из абзацев и таблиц документа для анализа агентами.
    Входной параметр file_path — путь к файлу."""
    if not os.path.exists(file_path):
        return f"Ошибка: файл '{file_path}' не найден."
    
    try:
        import docx
        doc = docx.Document(file_path)
        full_text = []
        
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                full_text.append(text)
                
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_text:
                    full_text.append(" | ".join(row_text))
                    
        if not full_text:
            return f"В Word-документе '{file_path}' не обнаружен текст."
            
        return "\n\n".join(full_text)
    except Exception as e:
        return f"Не удалось прочитать Word-документ: {e}"


# =====================================================================
# ИНСТРУМЕНТ 6. ПОИСК В БАЗЕ ДАННЫХ PUBMED (БИОМЕДИЦИНА, МЕДИЦИНА, ЭКОЛОГИЯ)
# =====================================================================
@tool
def search_pubmed(query: str) -> str:
    """Ищет научные медицинские, биологические и экологические публикации в базе PubMed.
    Возвращает список из 3 самых свежих статей с их названиями, авторами и источниками.
    Параметр query (поисковый запрос) должен быть на АНГЛИЙСКОМ языке (например, 'urban pollution air health')."""
    formatted_query = urllib.parse.quote(query.strip())
    search_url = f"https://eutils.ncbi.nlm.nih.gov/entrez/eutils/esearch.fcgi?db=pubmed&term={formatted_query}&retmode=json&retmax=3"
    headers = {'User-Agent': 'Mozilla/5.0'}
    
    try:
        req = urllib.request.Request(search_url, headers=headers)
        with urllib.request.urlopen(req, timeout=5) as res:
            data = json.loads(res.read())
        
        id_list = data.get('esearchresult', {}).get('idlist', [])
        if not id_list:
            return f"По запросу '{query}' в базе PubMed статей не найдено."
            
        ids = ",".join(id_list)
        summary_url = f"https://eutils.ncbi.nlm.nih.gov/entrez/eutils/esummary.fcgi?db=pubmed&id={ids}&retmode=json"
        
        req_sum = urllib.request.Request(summary_url, headers=headers)
        with urllib.request.urlopen(req_sum, timeout=5) as res_sum:
            sum_data = json.loads(res_sum.read())
            
        results = []
        uid_results = sum_data.get('result', {})
        for uid in id_list:
            article_info = uid_results.get(uid, {})
            if not article_info:
                continue
            title = article_info.get('title', 'Без названия')
            pub_date = article_info.get('pubdate', 'Неизвестно')
            authors = [a.get('name', '') for a in article_info.get('authors', [])]
            source = article_info.get('source', '')
            
            results.append(
                f"Название: {title}\n"
                f"Авторы: {', '.join(authors)}\n"
                f"Дата: {pub_date} | Источник: {source}\n"
                f"PubMed ID: {uid}\n"
                f"--------------------------------------------------"
            )
        
        return "\n".join(results) if results else "Не удалось получить детали для найденных статей."
        
    except Exception as e:
        return f"Ошибка при поиске на PubMed: {e}"


# =====================================================================
# ИНСТРУМЕНТ 7. ПОИСК В БАЗЕ ДАННЫХ OPENALEX (ОБЩЕНАУЧНЫЙ ПОИСК, 250 МЛН+ СТАТЕЙ)
# =====================================================================
@tool
def search_openalex(query: str) -> str:
    """Ищет междисциплинарные научные публикации по всему миру в открытой базе OpenAlex.
    Подходит для социально-экономических, гуманитарных, технологических и управленческих тем.
    Возвращает названия, авторов, даты публикаций и краткие резюме работ.
    Параметр query (поисковый запрос) должен быть на АНГЛИЙСКОМ языке (например, 'smart city traffic regulation')."""
    formatted_query = urllib.parse.quote(query.strip())
    url = f"https://api.openalex.org/works?search={formatted_query}&per_page=3"
    headers = {'User-Agent': 'Mozilla/5.0 (mailto:student@ranepa.ru)'}
    
    try:
        req = urllib.request.Request(url, headers=headers)
        with urllib.request.urlopen(req, timeout=5) as res:
            data = json.loads(res.read())
            
        results = data.get('results', [])
        if not results:
            return f"По запросу '{query}' в базе OpenAlex научных работ не найдено."
            
        output = []
        for work in results:
            title = work.get('title', 'Без названия')
            pub_year = work.get('publication_year', 'Не указан')
            
            authors_list = []
            for membership in work.get('authorships', []):
                author_name = membership.get('author', {}).get('display_name')
                if author_name:
                    authors_list.append(author_name)
                    
            abstract_index = work.get('abstract_inverted_index')
            abstract = "Аннотация отсутствует."
            if abstract_index:
                word_positions = {}
                for word, positions in abstract_index.items():
                    for pos in positions:
                        word_positions[pos] = word
                sorted_words = [word_positions[p] for p in sorted(word_positions.keys())]
                abstract = " ".join(sorted_words)
                if len(abstract) > 350:
                    abstract = abstract[:350] + "..."
                    
            output.append(
                f"Название: {title}\n"
                f"Авторы: {', '.join(authors_list[:5])}\n"
                f"Год: {pub_year}\n"
                f"Аннотация: {abstract}\n"
                f"--------------------------------------------------"
            )
            
        return "\n".join(output)
        
    except Exception as e:
        return f"Ошибка при поиске на OpenAlex: {e}"


# =====================================================================
# ИНСТРУМЕНТ 8. ЧТЕНИЕ ПРЕЗЕНТАЦИЙ POWERPOINT (.PPTX / .PPT)
# =====================================================================
@tool
def read_powerpoint_file(file_path: str) -> str:
    """Считывает текстовое содержимое презентации Microsoft PowerPoint (.pptx или .ppt).
    Извлекает текст со слайдов и текстовых блоков для анализа агентами.
    Входной параметр file_path — путь к файлу презентации."""
    if not os.path.exists(file_path):
        return f"Ошибка: файл '{file_path}' не найден."
    
    file_lower = file_path.lower()
    
    # Сценарий для PPTX
    if file_lower.endswith('.pptx'):
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            full_text = []
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                slide_text.append(text)
                if slide_text:
                    full_text.append(f"--- Слайд {i+1} ---\n" + "\n".join(slide_text))
            
            if not full_text:
                return f"В презентации PPTX '{file_path}' не обнаружен текст."
            return "\n\n".join(full_text)
        except Exception as e:
            return f"Не удалось прочитать презентацию .pptx: {e}"
            
    # Сценарий для PPT (бинарный формат)
    elif file_lower.endswith('.ppt'):
        try:
            with open(file_path, "rb") as f:
                data = f.read()
            
            import re
            # UTF-16LE строки длиной от 4 символов
            utf16_strings = re.findall(rb'(?:[\x20-\x7E\x00-\xFF]\x00){4,}', data)
            texts = []
            for s in utf16_strings:
                try:
                    txt = s.decode('utf-16le').strip()
                    if len(txt) > 3 and not any(junk in txt for junk in ['Component', 'Root', 'SummaryInformation', 'DocumentSummaryInformation', 'PowerPoint Document']):
                        texts.append(txt)
                except Exception:
                    pass
            
            # ASCII строки длиной от 4 символов
            ascii_strings = re.findall(rb'[\x20-\x7E]{4,}', data)
            for s in ascii_strings:
                try:
                    txt = s.decode('ascii').strip()
                    if len(txt) > 3 and not any(junk in txt for junk in ['Component', 'Root', 'SummaryInformation', 'DocumentSummaryInformation', 'PowerPoint Document']):
                        texts.append(txt)
                except Exception:
                    pass
            
            if not texts:
                return f"Не удалось автоматически извлечь текст из бинарного файла PPT '{file_path}'. Для надежности рекомендуется пересохранить файл в формате PPTX."
                
            unique_texts = sorted(list(set(texts)), key=lambda x: len(x), reverse=True)
            return (
                f"--- Извлеченный текст из бинарного файла PPT ({file_path}) ---\n"
                "(Внимание: текст получен прямым парсингом бинарных строк OLE, структура слайдов не сохранена)\n\n"
                + "\n\n".join(unique_texts[:100])
            )
        except Exception as e:
            return f"Ошибка при парсинге .ppt файла: {e}. Пожалуйста, пересохраните презентацию в формате .pptx."
            
    else:
        return f"Ошибка: неподдерживаемый формат презентации для файла '{file_path}'. Допустимы только .pptx и .ppt."
