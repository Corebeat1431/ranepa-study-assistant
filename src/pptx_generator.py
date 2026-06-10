import os
import json
import shutil
import re
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# =====================================================================
# ООП И МАКЕТИРОВАНИЕ ПРЕЗЕНТАЦИЙ (ДЛЯ ОБУЧЕНИЯ):
# 1. Объектная модель python-pptx:
#    - Presentation — это класс, представляющий весь файл презентации. 
#      Его экземпляр (объект prs) хранит коллекцию слайдов (prs.slides).
#    - Каждая фигура на слайде (карточка, текст, линия, картинка) является
#      объектом класса Shape. Мы добавляем их методами вроде slide.shapes.add_shape()
#      или slide.shapes.add_textbox() и настраиваем их свойства (fill, line, text_frame).
# 2. Адаптивность и масштабирование:
#    - Слайды верстаются в дюймах (Inches). Чтобы верстка хорошо смотрелась
#      как на экранах 4:3, так и 16:9, мы используем коэффициенты масштабирования scale_x и scale_y.
# =====================================================================
# Цветовая гамма РАНХиГС
MAROON = RGBColor(139, 0, 0)      # #8B0000 (основной бордовый)
DARK_GRAY = RGBColor(40, 40, 40)   # #282828 (цвет текста)
LIGHT_GRAY = RGBColor(128, 128, 128) # Цвет подзаголовков
FONT_NAME = "Montserrat"

def split_bullet(bullet_text):
    """Разделяет текст буллета на заголовок и описание по двоеточию, тире или первым 3 словам."""
    bullet_text = bullet_text.strip().lstrip("•-* ").strip()
    if ":" in bullet_text:
        parts = bullet_text.split(":", 1)
        return parts[0].strip(), parts[1].strip()
    elif "—" in bullet_text:
        parts = bullet_text.split("—", 1)
        return parts[0].strip(), parts[1].strip()
    else:
        words = bullet_text.split()
        if len(words) > 3:
            title = " ".join(words[:3])
            body = " ".join(words[3:])
            return title, body
        return bullet_text, ""

def extract_metric(bullet_text):
    """Находит первое число/процент/сумму в тексте и разделяет буллет на метрику и описание."""
    bullet_text = bullet_text.strip().lstrip("•-* ").strip()
    # Регулярное выражение ищет числа с суффиксами +, %, млн, тыс, руб, ₽
    match = re.search(r'(\d+[\s\-\+\%]*(?:тыс|млн|млрд|руб|₽|%|\+)?|\+\d+%)', bullet_text)
    if match:
        metric = match.group(1)
        description = bullet_text.replace(metric, "", 1).strip()
        # Очищаем стык от знаков препинания
        description = re.sub(r'^[\s\-\—\.\,\:]+', '', description).strip()
        return metric, description
    
    words = bullet_text.split()
    if len(words) > 2:
        return " ".join(words[:2]), " ".join(words[2:])
    return "Показатель", bullet_text

def translate_to_russian_via_gemini(text: str) -> str:
    """Проверяет наличие английских слов в тексте и переводит его на русский язык через Gemini API.
    
    Для обучения:
    - re.search(r'[a-zA-Z]') — регулярное выражение для поиска латинских букв в тексте.
    - В ООП мы используем genai.Client() для создания объекта-клиента и выполнения API-методов.
    """
    import re
    import os
    from dotenv import load_dotenv
    load_dotenv() # Загружаем ключи из .env для корректной авторизации API
    if not text or not re.search(r'[a-zA-Z]', text):
        return text # Перевод не требуется, если латиницы нет
        
    try:
        from google import genai
        api_key = os.getenv("GEMINI_API_KEY")
        if not api_key:
            return text
        client = genai.Client(api_key=api_key)
        # Делаем короткий вызов к самой быстрой модели
        response = client.models.generate_content(
            model="gemini-flash-lite-latest",
            contents=f"Translate this short text/label to Russian. Absolutely all terms, including IT, CAPEX, OPEX, KPI, and technical words, MUST be fully translated to Russian (e.g. CAPEX -> Капитальные затраты, IT -> ИТ, KPI -> КПЭ). Output ONLY the translated Russian text without any quotes, details, or explanations: '{text}'"
        )
        translated = response.text.strip().strip("'\"")
        print(f"[ПЕРЕВОД ГРАФИКА] Переведено: '{text}' -> '{translated}'")
        return translated

    except Exception as e:
        print(f"[ПЕРЕВОД ГРАФИКА] Ошибка при переводе '{text}': {e}")
        return text

def generate_chart_image(chart_data_dict):
    """Генерирует картинку графика с помощью matplotlib и сохраняет в output/"""
    if not chart_data_dict:
        return
    
    chart_type = chart_data_dict.get("chart_type", "bar")
    title = translate_to_russian_via_gemini(chart_data_dict.get("title", ""))
    labels = [translate_to_russian_via_gemini(lbl) for lbl in chart_data_dict.get("labels", [])]
    values = chart_data_dict.get("values", [])
    filename = chart_data_dict.get("filename", "")
    
    if not filename or not labels or not values:
        return
        
    try:
        # Цвета РАНХиГС
        primary_color = '#8B0000' # Темно-красный бордовый
        secondary_colors = ['#8B0000', '#A52A2A', '#D2691E', '#7F7F7F', '#A9A9A9', '#D3D3D3']
        
        plt.figure(figsize=(8, 5))
        plt.title(title, fontsize=14, fontweight='bold', pad=15, color='#333333')
        
        if chart_type.lower() == 'bar':
            bars = plt.bar(labels, values, color=primary_color, edgecolor='#555555', width=0.6)
            plt.grid(axis='y', linestyle='--', alpha=0.5)
            # Добавляем значения над столбцами
            for bar in bars:
                height = bar.get_height()
                plt.text(bar.get_x() + bar.get_width()/2.0, height, f'{height:g}', 
                         ha='center', va='bottom', fontsize=10, fontweight='bold')
                
        elif chart_type.lower() == 'line':
            plt.plot(labels, values, marker='o', linewidth=3, color=primary_color, markersize=8)
            plt.grid(True, linestyle='--', alpha=0.5)
            # Добавляем значения над точками
            for x, y in zip(labels, values):
                plt.text(x, y, f'{y:g}', ha='center', va='bottom', fontsize=10, fontweight='bold')
                
        elif chart_type.lower() == 'pie':
            plt.pie(values, labels=labels, autopct='%1.1f%%', 
                    colors=secondary_colors[:len(values)], startangle=140, 
                    textprops={'fontsize': 11, 'weight': 'bold'})
            
        else:
            print(f"[ГЕНЕРАТОР] Неподдерживаемый тип графика: {chart_type}")
            plt.close()
            return
        
        plt.tight_layout()
        os.makedirs("output", exist_ok=True)
        save_path = os.path.join("output", filename)
        plt.savefig(save_path, dpi=150)
        plt.close()
        print(f"[ГЕНЕРАТОР] График успешно создан в коде и сохранен: {save_path}")
    except Exception as e:
        plt.close()
        print(f"[ГЕНЕРАТОР] Ошибка при создании графика {filename}: {e}")

# Дополнительные функции верстки для повышения премиальности
def build_slide_2_cards(slide, bullets, scale_x, scale_y):
    """Строит сетку 2x2 модульных карточек на Слайде 2."""
    local_bullets = bullets[:4]
    card_coords = [
        (0.8, 2.0),
        (8.3, 2.0),
        (0.8, 5.2),
        (8.3, 5.2)
    ]
    
    for idx, bullet in enumerate(local_bullets):
        if idx >= len(card_coords):
            break
        x_val, y_val = card_coords[idx]
        x = x_val * scale_x
        y = y_val * scale_y
        w = 6.9 * scale_x
        h = 2.8 * scale_y
        
        # Добавляем карточку (скругленный прямоугольник)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(w), Inches(h)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(250, 245, 245) # Нежная теплая серо-розовая подложка
        card.line.color.rgb = RGBColor(220, 210, 210)
        card.line.width = Pt(1.5)
        
        # Тонкая декоративная полоса сверху карточки (РАНХиГС бордовый)
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y), Inches(w), Inches(0.08 * scale_y)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = MAROON
        accent_bar.line.color.rgb = MAROON
        
        # Разделяем буллеты на заголовок и тело
        b_title, b_body = split_bullet(bullet)
        
        # Добавляем текст
        tb = slide.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.25 * scale_y), Inches(w - 0.6), Inches(h - 0.45 * scale_y))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0)
        tf.margin_right = Inches(0)
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)
        
        p1 = tf.paragraphs[0]
        p1.text = f"0{idx+1}. {b_title}"
        p1.font.name = FONT_NAME
        p1.font.size = Pt(20)
        p1.font.bold = True
        p1.font.color.rgb = MAROON
        p1.space_after = Pt(6)
        
        if b_body:
            p2 = tf.add_paragraph()
            p2.text = b_body
            p2.font.name = FONT_NAME
            p2.font.size = Pt(15)
            p2.font.color.rgb = DARK_GRAY

def build_slide_8_metrics(slide, bullets, scale_x, scale_y):
    """Строит вертикальный список крупных числовых показателей на Слайде 8."""
    local_bullets = bullets[:4]
    row_h = 1.35 * scale_y
    start_y = 2.0 * scale_y
    
    for idx, bullet in enumerate(local_bullets):
        y = start_y + idx * (row_h + 0.25 * scale_y)
        
        # Выделяем цифру/метрику и описание
        metric, desc = extract_metric(bullet)
        
        # Текстовое поле для гигантской цифры
        tb_num = slide.shapes.add_textbox(Inches(0.8 * scale_x), Inches(y), Inches(2.2 * scale_x), row_h)
        tf_num = tb_num.text_frame
        tf_num.word_wrap = True
        p_num = tf_num.paragraphs[0]
        p_num.text = metric
        p_num.font.name = FONT_NAME
        p_num.font.size = Pt(40)
        p_num.font.bold = True
        p_num.font.color.rgb = MAROON
        
        # Текстовое поле для описания
        tb_desc = slide.shapes.add_textbox(Inches(3.1 * scale_x), Inches(y + 0.1 * scale_y), Inches(3.8 * scale_x), row_h)
        tf_desc = tb_desc.text_frame
        tf_desc.word_wrap = True
        p_desc = tf_desc.paragraphs[0]
        p_desc.text = desc
        p_desc.font.name = FONT_NAME
        p_desc.font.size = Pt(16)
        p_desc.font.color.rgb = DARK_GRAY
        
        # Рисуем разделительную линию под каждым показателем (кроме последнего)
        if idx < len(local_bullets) - 1:
            line_y = y + row_h + 0.1 * scale_y
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.8 * scale_x), Inches(line_y), Inches(6.1 * scale_x), Inches(0.02 * scale_y)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(220, 210, 210)
            line.line.color.rgb = RGBColor(220, 210, 210)

def is_stats_slide(bullets):
    """Определяет, содержит ли слайд количественные показатели/метрики."""
    count = 0
    for bullet in bullets:
        bullet_text = bullet.strip().lstrip("•-* ").strip()
        # Ищем числа или проценты в начале или в составе слов
        match = re.search(r'(\d+[\s\-\+\%]*(?:тыс|млн|млрд|руб|₽|%|\+)?|\+\d+%)', bullet_text)
        if match:
            count += 1
    return count >= 2

def draw_metric_card(slide, bullet, x, y, w, h, scale_y):
    """Рисует премиальную карточку метрики с большим числом сверху и описанием снизу."""
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(250, 245, 245)
    card.line.color.rgb = RGBColor(220, 210, 210)
    card.line.width = Pt(1.5)
    
    # Декоративный бордовый штрих сверху
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(0.08 * scale_y)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = MAROON
    accent_bar.line.color.rgb = MAROON
    
    metric, desc = extract_metric(bullet)
    
    # Числовой показатель
    tb_num = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.25 * scale_y), Inches(w - 0.4), Inches(1.3 * scale_y))
    tf_num = tb_num.text_frame
    tf_num.word_wrap = True
    tf_num.margin_left = Inches(0)
    tf_num.margin_right = Inches(0)
    tf_num.margin_top = Inches(0)
    tf_num.margin_bottom = Inches(0)
    
    p_num = tf_num.paragraphs[0]
    p_num.text = metric
    p_num.font.name = FONT_NAME
    p_num.font.size = Pt(40)
    p_num.font.bold = True
    p_num.font.color.rgb = MAROON
    
    # Описание показателя
    tb_desc = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 1.6 * scale_y), Inches(w - 0.4), Inches(h - 1.8 * scale_y))
    tf_desc = tb_desc.text_frame
    tf_desc.word_wrap = True
    tf_desc.margin_left = Inches(0)
    tf_desc.margin_right = Inches(0)
    tf_desc.margin_top = Inches(0)
    tf_desc.margin_bottom = Inches(0)
    
    p_desc = tf_desc.paragraphs[0]
    p_desc.text = desc
    p_desc.font.name = FONT_NAME
    p_desc.font.size = Pt(13)
    p_desc.font.color.rgb = DARK_GRAY

def render_title_slide(slide, title_text, use_template, scale_x, scale_y):
    """Отрисовка титульного слайда (поддерживает шаблон и создание с нуля)."""
    if use_template:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                if "Тема выпускной" in text or "индивидуальный проект" in text or "Тема проекта" in text:
                    tf = shape.text_frame
                    tf.clear()
                    p = tf.paragraphs[0]
                    p.text = title_text.upper()
                    p.font.name = FONT_NAME
                    p.font.size = Pt(32)
                    p.font.bold = True
                    p.font.color.rgb = MAROON
                elif "Слушатель программы" in text or "Руководитель" in text or "Наставник" in text:
                    tf = shape.text_frame
                    tf.clear()
                    p = tf.paragraphs[0]
                    p.text = (
                        "Слушатель программы: _______________________\n"
                        "Руководитель проектной работы: _______________________"
                    )
                    p.font.name = FONT_NAME
                    p.font.size = Pt(18)
                    p.font.color.rgb = DARK_GRAY
    else:
        # Если создаем с нуля
        title_box = slide.shapes.add_textbox(Inches(1.0 * scale_x), Inches(2.0 * scale_y), Inches(14.0 * scale_x), Inches(2.5 * scale_y))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text.upper()
        p.font.name = FONT_NAME
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = MAROON
        
        sub_box = slide.shapes.add_textbox(Inches(1.0 * scale_x), Inches(4.5 * scale_y), Inches(14.0 * scale_x), Inches(2.0 * scale_y))
        tf_sub = sub_box.text_frame
        tf_sub.word_wrap = True
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = (
            "Слушатель программы: _______________________\n"
            "Руководитель проектной работы: _______________________"
        )
        p_sub.font.name = FONT_NAME
        p_sub.font.size = Pt(18)
        p_sub.font.color.rgb = DARK_GRAY

def sanitize_image_format(image_path):
    """Проверяет формат изображения с помощью Pillow.
    
    Если изображение имеет формат, не поддерживаемый python-pptx (например, MPO или WEBP),
    конвертирует его в стандартный JPEG, чтобы избежать падения при сборке слайда.
    """
    from PIL import Image
    try:
        # Проверяем формат файла
        with Image.open(image_path) as img:
            fmt = img.format
            # Поддерживаемые форматы python-pptx: BMP, GIF, JPEG, PNG, TIFF, WMF
            if fmt not in ['BMP', 'GIF', 'JPEG', 'PNG', 'TIFF', 'WMF']:
                print(f"[САНИТАЙЗЕР ИЗОБРАЖЕНИЙ] Формат '{fmt}' для '{image_path}' не поддерживается. Конвертируем в JPEG...")
                rgb_img = img.convert('RGB')
                rgb_img.save(image_path, 'JPEG')
                print("[САНИТАЙЗЕР ИЗОБРАЖЕНИЙ] Конвертация завершена успешно.")
    except Exception as e:
        print(f"[САНИТАЙЗЕР ИЗОБРАЖЕНИЙ] Ошибка при проверке или конвертации формата файла '{image_path}': {e}")

def render_split_slide(slide, bullets, image_name, scale_x, scale_y):
    """Сплит-слайд: левая колонка - структурированный текст, правая - изображение/график."""
    x = 0.8 * scale_x
    y = 2.0 * scale_y
    w = 6.0 * scale_x
    h = 6.0 * scale_y
    
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    
    # Динамически распределяем размер текста и отступы в зависимости от числа тезисов
    if len(bullets) <= 2:
        font_size = Pt(19)
        spacing = Pt(22)
    elif len(bullets) == 3:
        font_size = Pt(17)
        spacing = Pt(16)
    else:
        font_size = Pt(15)
        spacing = Pt(12)

    for k, bullet in enumerate(bullets):
        p = tf.add_paragraph() if k > 0 else tf.paragraphs[0]
        clean_bullet = bullet.strip().lstrip("•-* ").strip()
        p.text = "→   " + clean_bullet
        p.space_after = spacing
        p.font.name = FONT_NAME
        p.font.size = font_size
        p.font.color.rgb = DARK_GRAY
        
        #/ Красиво выделяем заголовки буллетов жирным бордовым цветом
        if ":" in clean_bullet:
            parts = clean_bullet.split(":", 1)
            p.text = "→   "
            run_bold = p.add_run()
            run_bold.text = parts[0].strip() + ": "
            run_bold.font.bold = True
            run_bold.font.color.rgb = MAROON
            run_bold.font.name = FONT_NAME
            run_bold.font.size = font_size
            
            run_desc = p.add_run()
            run_desc.text = parts[1].strip()
            run_desc.font.color.rgb = DARK_GRAY
            run_desc.font.name = FONT_NAME
            run_desc.font.size = font_size
            
    # Рисуем изображение в правой части
    resolved_path = ""
    if image_name:
        for path_candidate in [os.path.join("output", image_name), os.path.join("assets", "images", image_name), image_name]:
            if os.path.isfile(path_candidate):
                resolved_path = path_candidate
                break
            
    if resolved_path:
        sanitize_image_format(resolved_path)
        slide.shapes.add_picture(
            resolved_path, 
            Inches(7.5 * scale_x), 
            Inches(2.0 * scale_y), 
            Inches(7.5 * scale_x), 
            Inches(5.5 * scale_y)
        )
    else:
        print(f"[ПРЕДУПРЕЖДЕНИЕ] Изображение '{image_name}' не найдено ни в output, ни в assets/images")

def render_stats_layout(slide, bullets, scale_x, scale_y):
    """Верстка крупных показателей в виде горизонтального ряда карточек."""
    local_bullets = bullets[:4]
    num_items = len(local_bullets)
    
    if num_items == 3:
        w_card = 4.4 * scale_x
        gap = 0.6 * scale_x
        start_x = 0.8 * scale_x
        y = 2.5 * scale_y
        h = 4.5 * scale_y
        for idx, bullet in enumerate(local_bullets):
            draw_metric_card(slide, bullet, start_x + idx * (w_card + gap), y, w_card, h, scale_y)
            
    elif num_items == 4:
        w_card = 3.25 * scale_x
        gap = 0.5 * scale_x
        start_x = 0.8 * scale_x
        y = 2.5 * scale_y
        h = 4.5 * scale_y
        for idx, bullet in enumerate(local_bullets):
            draw_metric_card(slide, bullet, start_x + idx * (w_card + gap), y, w_card, h, scale_y)
            
    else:
        # Вертикальный список крупных цифр, если элементов другое число
        build_slide_8_metrics(slide, bullets, scale_x, scale_y)

def render_three_columns_layout(slide, bullets, scale_x, scale_y):
    """Верстка 3-х колонок карточек для слайда с тремя тезисами."""
    w_card = 4.4 * scale_x
    gap = 0.6 * scale_x
    start_x = 0.8 * scale_x
    y = 2.2 * scale_y
    h = 5.5 * scale_y
    
    for idx, bullet in enumerate(bullets):
        x = start_x + idx * (w_card + gap)
        
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(w_card), Inches(h)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(250, 245, 245)
        card.line.color.rgb = RGBColor(220, 210, 210)
        card.line.width = Pt(1.5)
        
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y), Inches(w_card), Inches(0.08 * scale_y)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = MAROON
        accent_bar.line.color.rgb = MAROON
        
        b_title, b_body = split_bullet(bullet)
        
        tb = slide.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.3 * scale_y), Inches(w_card - 0.6), Inches(h - 0.5 * scale_y))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0)
        tf.margin_right = Inches(0)
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)
        
        p1 = tf.paragraphs[0]
        p1.text = f"0{idx+1}. {b_title}"
        p1.font.name = FONT_NAME
        p1.font.size = Pt(20)
        p1.font.bold = True
        p1.font.color.rgb = MAROON
        p1.space_after = Pt(8)
        
        if b_body:
            p2 = tf.add_paragraph()
            p2.text = b_body
            p2.font.name = FONT_NAME
            p2.font.size = Pt(14)
            p2.font.color.rgb = DARK_GRAY

def render_default_bullet_layout(slide, bullets, scale_x, scale_y):
    """Дефолтная верстка списком во всю ширину слайда."""
    x = 1.0 * scale_x
    y = 2.0 * scale_y
    w = 14.0 * scale_x
    h = 6.0 * scale_y
    
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    
    # Динамически распределяем отступы и размер текста для равномерного заполнения слайда
    if len(bullets) <= 2:
        font_size = Pt(22)
        spacing = Pt(28)
    elif len(bullets) == 3:
        font_size = Pt(20)
        spacing = Pt(22)
    else:
        font_size = Pt(18)
        spacing = Pt(16)

    for k, bullet in enumerate(bullets):
        p = tf.add_paragraph() if k > 0 else tf.paragraphs[0]
        clean_bullet = bullet.strip().lstrip("•-* ").strip()
        p.text = "→   " + clean_bullet
        p.space_after = spacing
        p.font.name = FONT_NAME
        p.font.size = font_size
        p.font.color.rgb = DARK_GRAY
        
        if ":" in clean_bullet:
            parts = clean_bullet.split(":", 1)
            p.text = "→   "
            run_bold = p.add_run()
            run_bold.text = parts[0].strip() + ": "
            run_bold.font.bold = True
            run_bold.font.color.rgb = MAROON
            run_bold.font.name = FONT_NAME
            run_bold.font.size = font_size
            
            run_desc = p.add_run()
            run_desc.text = parts[1].strip()
            run_desc.font.color.rgb = DARK_GRAY
            run_desc.font.name = FONT_NAME
            run_desc.font.size = font_size

def generate_fallback_chart(filename, bullets):
    """Генерирует резервный график, если файл не был создан агентом."""
    import matplotlib.pyplot as plt
    try:
        plt.figure(figsize=(8, 5))
        primary_color = '#8B0000' # Темно-красный бордовый
        
        if "4" in filename:
            title = "SWOT-оценка факторов проекта"
            labels = ["Сильные\nстороны", "Слабости", "Возможности", "Угрозы"]
            values = [8.5, 4.0, 9.0, 5.5]
            
            bars = plt.bar(labels, values, color=['#8B0000', '#7F7F7F', '#A52A2A', '#A9A9A9'], edgecolor='#555555', width=0.5)
            plt.ylabel("Влияние (баллы)", fontsize=11, fontweight='bold')
            plt.grid(axis='y', linestyle='--', alpha=0.5)
            plt.ylim(0, 10)
            for bar in bars:
                height = bar.get_height()
                plt.text(bar.get_x() + bar.get_width()/2.0, height + 0.2, f'{height:g}', 
                         ha='center', va='bottom', fontsize=10, fontweight='bold')
                         
        elif "7" in filename:
            title = "Прогноз финансовых показателей (млн руб.)"
            labels = ["Капвложения", "Выручка (Г1)", "Выручка (Г2)", "Выручка (Г3)"]
            values = [15.0, 6.0, 14.5, 22.0]
            
            plt.plot(labels, values, marker='o', linewidth=3, color=primary_color, markersize=8)

            plt.ylabel("Млн рублей", fontsize=11, fontweight='bold')
            plt.grid(True, linestyle='--', alpha=0.5)
            for x, y in zip(labels, values):
                plt.text(x, y, f'{y:g}', ha='center', va='bottom', fontsize=10, fontweight='bold')
        else:
            plt.close()
            return
            
        plt.title(title, fontsize=13, fontweight='bold', pad=15, color='#333333')
        plt.tight_layout()
        os.makedirs("output", exist_ok=True)
        save_path = os.path.join("output", filename)
        plt.savefig(save_path, dpi=150)
        plt.close()
        print(f"[РЕЗЕРВ] Сгенерирован график-заглушка: {save_path}")
    except Exception as e:
        plt.close()
        print(f"[РЕЗЕРВ] Ошибка генерации графика-заглушки: {e}")

def get_keywords_from_title(title):
    """Возвращает английские ключевые слова для поиска картинок по заголовку слайда."""
    title_lower = title.lower()
    if "актуальн" in title_lower or "тренд" in title_lower:
        return "tourism travel"
    if "наук" in title_lower or "исслед" in title_lower or "институт" in title_lower or "потенциал" in title_lower:
        return "science research lab"
    if "инфраструктур" in title_lower or "цифр" in title_lower or "решен" in title_lower:
        return "technology digital innovation"
    if "сезон" in title_lower or "школьн" in title_lower or "аудитор" in title_lower:
        return "education classroom student"
    if "социальн" in title_lower or "эффект" in title_lower or "рабоч" in title_lower or "ваканс" in title_lower:
        return "team success city job"
    if "дорожн" in title_lower or "карт" in title_lower or "этап" in title_lower or "реализ" in title_lower:
        return "planning roadmap timeline"
    if "вывод" in title_lower or "заключен" in title_lower:
        return "city landscape nature"
    return "city business work"

def get_thematic_local_backup(keyword: str) -> str:
    """Определяет наиболее подходящее имя локального резервного изображения по ключевым словам.
    
    Для обучения:
    - Метод .lower() приводит строку к нижнему регистру для регистронезависимого сравнения.
    - Оператор 'in' проверяет вхождение подстроки в строку.
    """
    kw_lower = keyword.lower()
    if "sevastopol" in kw_lower or "севастополь" in kw_lower:
        return "sevastopol.jpg"
    if "science" in kw_lower or "research" in kw_lower or "наук" in kw_lower or "исслед" in kw_lower:
        return "science.jpg"
    if "tourism" in kw_lower or "travel" in kw_lower or "туриз" in kw_lower:
        return "tourism.jpg"
    if "tech" in kw_lower or "digital" in kw_lower or "цифр" in kw_lower or "инновац" in kw_lower:
        return "technology.jpg"
    if "edu" in kw_lower or "student" in kw_lower or "обуч" in kw_lower or "молодеж" in kw_lower:
        return "education.jpg"
    if "team" in kw_lower or "social" in kw_lower or "work" in kw_lower or "социаль" in kw_lower or "youth" in kw_lower or "community" in kw_lower:
        return "team.jpg"
    if "plan" in kw_lower or "budget" in kw_lower or "risk" in kw_lower or "карт" in kw_lower or "бюджет" in kw_lower or "риск" in kw_lower or "profit" in kw_lower or "revenue" in kw_lower:
        return "business.jpg"
    return "sevastopol.jpg"

def download_thematic_image(keyword, filename):
    """Генерирует и скачивает тематическое изображение с помощью свободных API на Hugging Face.
    
    Для обеспечения высокого качества и соответствия теме («Модернизация Инвестиционного
    портала города Севастополя» или любой другой теме), мы динамически генерируем
    новые картинки с помощью нейросетей SDXL-Lightning и SD 3.5 Large через Gradio Client.
    Если генерация завершается сбоем, мы используем локальные резервные шаблоны.
    """
    import os
    import shutil
    import re
    from PIL import Image
    
    os.makedirs("output", exist_ok=True)
    target_path = os.path.join("output", filename)
    
    # Пул известных размеров локальных шаблонов, чтобы перезаписать их, если они уже скопированы
    template_sizes = {100113, 86980, 51474, 57568, 624793, 808391, 148051}
    
    if os.path.exists(target_path):
        file_size = os.path.getsize(target_path)
        if file_size in template_sizes:
            print(f"[ДИНАМИЧЕСКИЙ ПОДБОР] Обнаружен файл-шаблон ({file_size} байт). Удаляем для повторной генерации.")
            try:
                os.remove(target_path)
            except Exception as e:
                print(f"[ДИНАМИЧЕСКИЙ ПОДБОР] Не удалось удалить шаблон: {e}")
        else:
            print(f"[ДИНАМИЧЕСКИЙ ПОДБОР] Файл {target_path} уже сгенерирован ранее. Пропускаем.")
            return filename
            
    # Пытаемся извлечь тему из имени файла
    theme_part = ""
    slide_num = "1"
    match = re.search(r'slide_(\d+)_(.+)_theme\.jpg', filename)
    if match:
        slide_num = match.group(1)
        theme_part = match.group(2).replace("_", " ").strip()
    else:
        match2 = re.search(r'slide_(\d+)_', filename)
        if match2:
            slide_num = match2.group(1)
        
    print(f"[ДИНАМИЧЕСКИЙ ПОДБОР] Начало генерации для: тема='{theme_part}', ключевые слова='{keyword}'")
    
    try:
        from gradio_client import Client
        
        # Карта ключевых слов для детальных промптов
        kw_map = {
            "tourism travel": "tourism infrastructure, travel guide, city landmark",
            "science research lab": "scientific research workspace, data statistics charts, modern laboratory",
            "technology digital innovation": "digital interface dashboard, web portal technology, computer screen with financial charts",
            "education classroom student": "education and professional business training, modern lecture hall, presentation",
            "team success city job": "corporate business team working, professional cooperation, project discussion",
            "planning roadmap timeline": "business roadmap layout, strategic timeline visualization, growth steps",
            "city landscape nature": "modern urban infrastructure, city skyline, development architecture",
            "city business work": "corporate workspace development, professional office environment, investment focus"
        }
        
        clean_kw = kw_map.get(keyword.strip().lower(), keyword)
        
        # Формируем промпт в стиле РАНХиГС с бордовыми акцентами
        prompt = (
            f"3D render, premium corporate business illustration, {clean_kw} concept, "
            f"themed around {theme_part if theme_part else 'business development and investment'}, "
            f"modern style, color palette of maroon, dark red, white and grey, "
            f"clean minimal composition, studio lighting, high resolution, "
            f"strictly no text, no letters, no words, no charts with text, no writing, no labels, no watermark"
        )
        print(f"[ДИНАМИЧЕСКИЙ ПОДБОР] Промпт: '{prompt}'")
        
        generated_img_path = None
        
        # Шаг 1. Пытаемся использовать ap123/SDXL-Lightning (самый быстрый)
        print(f"[ДИНАМИЧЕСКИЙ ПОДБОР] Попытка генерации через ap123/SDXL-Lightning...")
        try:
            client = Client("ap123/SDXL-Lightning")
            result_path = client.predict(
                prompt,
                "4-Step",
                api_name="/generate_image"
            )
            if result_path and os.path.exists(result_path):
                generated_img_path = result_path
                print(f"[ДИНАМИЧЕСКИЙ ПОДБОР] Успешно сгенерировано через ap123/SDXL-Lightning: {result_path}")
        except Exception as e_ap123:
            print(f"[ДИНАМИЧЕСКИЙ ПОДБОР] Сбой на ap123/SDXL-Lightning: {e_ap123}")
            
        # Шаг 2. Если не получилось, пытаемся использовать stabilityai/stable-diffusion-3.5-large
        if not generated_img_path:
            print(f"[ДИНАМИЧЕСКИЙ ПОДБОР] Попытка генерации через stabilityai/stable-diffusion-3.5-large...")
            try:
                client = Client("stabilityai/stable-diffusion-3.5-large")
                result = client.predict(
                    prompt,                                    # prompt
                    "ugly, deformed, text, word, label, sign", # negative_prompt
                    0,                                         # seed
                    True,                                      # randomize_seed
                    1024,                                      # width
                    1024,                                      # height
                    4.5,                                       # guidance_scale
                    20,                                        # num_inference_steps
                    api_name="/infer"
                )
                if result and isinstance(result, tuple) and len(result) > 0:
                    result_path = result[0]
                    if result_path and os.path.exists(result_path):
                        generated_img_path = result_path
                        print(f"[ДИНАМИЧЕСКИЙ ПОДБОР] Успешно сгенерировано через SD 3.5 Large: {result_path}")
            except Exception as e_sd35:
                print(f"[ДИНАМИЧЕСКИЙ ПОДБОР] Сбой на stabilityai/stable-diffusion-3.5-large: {e_sd35}")
                
        # Шаг 3. Если сгенерировали, копируем и проверяем
        if generated_img_path:
            shutil.copy2(generated_img_path, target_path)
            # Проверяем изображение
            with Image.open(target_path) as img:
                print(f"[ДИНАМИЧЕСКИЙ ПОДБОР] Изображение успешно получено и верифицировано: {img.format}, {img.size}")
            return filename
        else:
            raise RuntimeError("Все доступные Gradio Spaces не смогли сгенерировать изображение.")
            
    except Exception as e:
        print(f"[ДИНАМИЧЕСКИЙ ПОДБОР] Ошибка при генерации/загрузке изображения: {e}")
        print(f"[ДИНАМИЧЕСКИЙ ПОДБОР] Откатываемся к локальному шаблону.")
        
        # Сначала проверяем наличие уникального сгенерированного слайда
        custom_backup_png = f"slide_{slide_num}.png"
        custom_backup_jpg = f"slide_{slide_num}.jpg"
        
        local_backup = None
        for cb in [custom_backup_png, custom_backup_jpg]:
            if os.path.exists(os.path.join("assets", "images", cb)):
                local_backup = cb
                break
                
        if local_backup:
            ext = os.path.splitext(local_backup)[1]
            base_name = os.path.splitext(filename)[0]
            filename = base_name + ext
            target_path = os.path.join("output", filename)
        else:
            local_backup = get_thematic_local_backup(keyword)
            
        local_path = os.path.join("assets", "images", local_backup)
        if os.path.exists(local_path):
            try:
                shutil.copy2(local_path, target_path)
                print(f"[ДИНАМИЧЕСКИЙ ПОДБОР] Скопирован резервный файл {local_backup} -> {filename}")
                return filename
            except Exception as copy_err:
                print(f"[ДИНАМИЧЕСКИЙ ПОДБОР] Ошибка при копировании резерва: {copy_err}")
                
        return None


def draw_vector_timeline(slide, bullets, scale_x, scale_y):
    """Рисует красивую вертикальную дорожную карту (таймлайн) на правой стороне слайда.
    
    Для обучения:
    - MSO_SHAPE.OVAL — круглые маркеры для этапов дорожной карты.
    - MSO_SHAPE.RECTANGLE — тонкие полосы, используемые как линии связей.
    """
    x_offset = 7.5 * scale_x
    y_offset = 2.0 * scale_y
    w_area = 5.0 * scale_x
    h_area = 5.0 * scale_y
    
    local_bullets = bullets[:4] # Не более 4 этапов
    num_steps = len(local_bullets)
    if num_steps == 0:
        return
        
    step_gap = h_area / max(num_steps, 1)
    
    # 1. Рисуем вертикальную направляющую линию (ось времени)
    line_x = x_offset + 0.8 * scale_x
    line_y_start = y_offset + 0.3 * scale_y
    line_y_end = y_offset + (num_steps - 0.7) * step_gap
    
    time_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(line_x - 0.02 * scale_x), Inches(line_y_start), 
        Inches(0.04 * scale_x), Inches(line_y_end - line_y_start)
    )
    time_line.fill.solid()
    time_line.fill.fore_color.rgb = RGBColor(220, 210, 210)
    time_line.line.color.rgb = RGBColor(220, 210, 210)
    
    # 2. Рисуем шаги
    for idx, bullet in enumerate(local_bullets):
        cy = y_offset + 0.3 * scale_y + idx * step_gap
        
        # Рисуем круг-маркер
        r = 0.5 * scale_y
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(line_x - r/2), Inches(cy), Inches(r), Inches(r)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = MAROON
        circle.line.color.rgb = MAROON
        
        # Цифра внутри круга
        tf_c = circle.text_frame
        tf_c.clear()
        tf_c.margin_left = Inches(0)
        tf_c.margin_right = Inches(0)
        tf_c.margin_top = Inches(0)
        tf_c.margin_bottom = Inches(0)
        
        p_c = tf_c.paragraphs[0]
        p_c.text = str(idx + 1)
        p_c.alignment = 1 # Center
        p_c.font.name = FONT_NAME
        p_c.font.size = Pt(14)
        p_c.font.bold = True
        p_c.font.color.rgb = RGBColor(255, 255, 255)
        
        # Текст этапа справа от круга
        tb_w = w_area - 1.5 * scale_x
        tb = slide.shapes.add_textbox(
            Inches(line_x + r/2 + 0.3 * scale_x), Inches(cy - 0.05 * scale_y), 
            Inches(tb_w), Inches(step_gap - 0.1 * scale_y)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0)
        tf.margin_right = Inches(0)
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)
        
        b_title, b_body = split_bullet(bullet)
        
        p1 = tf.paragraphs[0]
        p1.text = b_title
        p1.font.name = FONT_NAME
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = MAROON
        
        if b_body:
            p2 = tf.add_paragraph()
            p2.text = b_body
            p2.font.name = FONT_NAME
            p2.font.size = Pt(13)
            p2.font.color.rgb = DARK_GRAY


def draw_vector_structure(slide, bullets, title_text, scale_x, scale_y):
    """Рисует структурную схему с центральным концептом и связями к тезисам.
    
    Для обучения:
    - Блок-схема строится по принципу иерархического дерева (Tree diagram).
    - Тонкие прямоугольники играют роль связующих линий между объектами.
    """
    x_offset = 7.2 * scale_x
    y_offset = 2.0 * scale_y
    w_area = 5.3 * scale_x
    h_area = 5.0 * scale_y
    
    local_bullets = bullets[:3] # Ограничимся 3 дочерними блоками для чистоты
    num_items = len(local_bullets)
    if num_items == 0:
        return
        
    # Координаты центрального блока
    cx = x_offset + 0.2 * scale_x
    cy = y_offset + 2.0 * scale_y
    cw = 2.2 * scale_x
    ch = 1.2 * scale_y
    
    # 1. Центральный блок (эллипс)
    center_shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(cx), Inches(cy), Inches(cw), Inches(ch)
    )
    center_shape.fill.solid()
    center_shape.fill.fore_color.rgb = MAROON
    center_shape.line.color.rgb = MAROON
    
    tf_c = center_shape.text_frame
    tf_c.word_wrap = True
    tf_c.margin_left = Inches(0.05)
    tf_c.margin_right = Inches(0.05)
    tf_c.margin_top = Inches(0.05)
    tf_c.margin_bottom = Inches(0.05)
    
    p_c = tf_c.paragraphs[0]
    words = title_text.split()
    p_c.text = " ".join(words[:2]).upper() if words else "ПРОЕКТ"
    p_c.alignment = 1 # Center
    p_c.font.name = FONT_NAME
    p_c.font.size = Pt(14)
    p_c.font.bold = True
    p_c.font.color.rgb = RGBColor(255, 255, 255)
    
    # Координаты дочерних блоков
    child_coords = []
    if num_items == 1:
        child_coords = [(cx + 2.8 * scale_x, cy)]
    elif num_items == 2:
        child_coords = [
            (cx + 2.8 * scale_x, cy - 1.2 * scale_y),
            (cx + 2.8 * scale_x, cy + 1.2 * scale_y)
        ]
    elif num_items >= 3:
        child_coords = [
            (cx + 2.8 * scale_x, cy - 1.8 * scale_y),
            (cx + 2.8 * scale_x, cy),
            (cx + 2.8 * scale_x, cy + 1.8 * scale_y)
        ]
        
    # 2. Отрисовка линий связи и дочерних блоков
    x_trunk = cx + cw + 0.3 * scale_x
    # Высоту линии развилки рассчитываем исходя из реальной высоты дочерних блоков th = 1.4 * scale_y
    y_start = min(ty + (1.4 * scale_y)/2 for tx, ty in child_coords[:num_items])
    y_end = max(ty + (1.4 * scale_y)/2 for tx, ty in child_coords[:num_items])
    
    # Горизонтальный ствол из центра к развилке
    stem = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(cx + cw), Inches(cy + ch/2 - 0.01 * scale_y),
        Inches(0.3 * scale_x), Inches(0.02 * scale_y)
    )
    stem.fill.solid()
    stem.fill.fore_color.rgb = RGBColor(220, 210, 210)
    stem.line.color.rgb = RGBColor(220, 210, 210)
    
    # Вертикальная линия развилки (если элементов больше 1)
    if num_items > 1:
        trunk = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x_trunk), Inches(y_start),
            Inches(0.02 * scale_x), Inches(y_end - y_start)
        )
        trunk.fill.solid()
        trunk.fill.fore_color.rgb = RGBColor(220, 210, 210)
        trunk.line.color.rgb = RGBColor(220, 210, 210)
        
    for idx, (tx, ty) in enumerate(child_coords[:num_items]):
        bullet = local_bullets[idx]
        b_title, b_body = split_bullet(bullet)
        
        tw = 2.5 * scale_x
        th = 1.4 * scale_y
        
        # Горизонтальная ветка к дочернему блоку
        branch_w = tx - x_trunk
        branch = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x_trunk), Inches(ty + th/2 - 0.01 * scale_y),
            Inches(branch_w), Inches(0.02 * scale_y)
        )
        branch.fill.solid()
        branch.fill.fore_color.rgb = RGBColor(220, 210, 210)
        branch.line.color.rgb = RGBColor(220, 210, 210)
        
        # Дочерний скругленный прямоугольник
        child_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(tx), Inches(ty), Inches(tw), Inches(th)
        )
        child_shape.fill.solid()
        child_shape.fill.fore_color.rgb = RGBColor(250, 245, 245)
        child_shape.line.color.rgb = MAROON
        child_shape.line.width = Pt(1.5)
        
        tf_t = child_shape.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left = Inches(0.08)
        tf_t.margin_right = Inches(0.08)
        tf_t.margin_top = Inches(0.08)
        tf_t.margin_bottom = Inches(0.08)
        
        p_t = tf_t.paragraphs[0]
        p_t.text = b_title
        p_t.font.name = FONT_NAME
        p_t.font.size = Pt(13)
        p_t.font.bold = True
        p_t.font.color.rgb = MAROON
        
        if b_body:
            p_t2 = tf_t.add_paragraph()
            body_words = b_body.split()
            # Берем до 6 слов, чтобы текст гарантированно помещался внутри цветного блока
            short_body = " ".join(body_words[:6]) + ("..." if len(body_words) > 6 else "")
            p_t2.text = short_body
            p_t2.font.name = FONT_NAME
            p_t2.font.size = Pt(10)
            p_t2.font.color.rgb = DARK_GRAY


def classify_and_render_slide(slide, title_text, bullets, image_name, i, scale_x, scale_y, use_template):
    """Классифицирует контент слайда и выбирает наилучшую верстку."""
    slide_num = i + 1
    
    # 1. Автоматический подбор иллюстраций из открытых источников для ритма презентации
    # Мы подбираем картинки так, чтобы чередовать карточки и сплит-слайды
    if not image_name:
        if slide_num == 4:
            image_name = "slide_4_chart.png"
        elif slide_num == 7:
            image_name = "slide_7_chart.png"
        elif slide_num in [3, 6, 9, 10]:
            # Вместо скачивания фотографий мы теперь используем векторные блок-схемы
            image_name = None

    # 2. Если это график и его нет на диске, создаем его
    if image_name in ["slide_4_chart.png", "slide_7_chart.png"]:
        img_path = os.path.join("output", image_name)
        if not os.path.exists(img_path):
            generate_fallback_chart(image_name, bullets)

    # 3. Проверяем существование файла картинки
    has_image = False
    if image_name:
        for path_candidate in [os.path.join("output", image_name), os.path.join("assets", "images", image_name), image_name]:
            if os.path.isfile(path_candidate):
                has_image = True
                break
                
    # СЛАЙД 1: Титульный
    if i == 0:
        render_title_slide(slide, title_text, use_template, scale_x, scale_y)
        return
        
    # Очищаем старые плейсхолдеры шаблона
    if use_template:
        text_shapes = [s for s in slide.shapes if s.has_text_frame]
        text_shapes.sort(key=lambda s: s.top if s.top else 0)
        if text_shapes:
            title_shape = text_shapes[0]
            tf_title = title_shape.text_frame
            tf_title.clear()
            p_title = tf_title.paragraphs[0]
            p_title.text = title_text
            p_title.font.name = FONT_NAME
            p_title.font.size = Pt(28)
            p_title.font.bold = True
            p_title.font.color.rgb = MAROON
            
            for extra in text_shapes[1:]:
                sp = extra._element
                sp.getparent().remove(sp)
    else:
        # Отрисовка заголовка с нуля
        header_box = slide.shapes.add_textbox(Inches(1.0 * scale_x), Inches(0.5 * scale_y), Inches(14.0 * scale_x), Inches(1.0 * scale_y))
        tf_header = header_box.text_frame
        tf_header.word_wrap = True
        p_head = tf_header.paragraphs[0]
        p_head.text = title_text
        p_head.font.name = FONT_NAME
        p_head.font.size = Pt(32)
        p_head.font.bold = True
        p_head.font.color.rgb = MAROON

    # Маршрутизируем контент по макетам
    # Приоритет макетов для создания профессионального визуального ритма РАНХиГС:
    # 1. Если это Слайд 6 -> Сплит-слайд с вертикальной дорожной картой (Timeline)
    # 2. Если это Слайд 3, 9 или 10 -> Сплит-слайд с векторной Mind Map структурой
    # 3. Если это график или принудительно заданная картинка (не авто-подбор), всегда сплит-слайд
    # 4. Если на слайде много числовых показателей -> крупные метрики
    # 5. Если ровно 3 тезиса -> 3 колонки
    # 6. Если ровно 4 тезиса -> 2x2 карточки
    # 7. Иначе -> дефолтный список
    
    is_auto_image = (image_name in ["tourism.jpg", "science.jpg", "technology.jpg", "education.jpg", "team.jpg", "business.jpg", "sevastopol.jpg"])
    
    if slide_num == 6:
        render_split_slide(slide, bullets, "", scale_x, scale_y)
        draw_vector_timeline(slide, bullets, scale_x, scale_y)
    elif slide_num in [3, 9, 10]:
        render_split_slide(slide, bullets, "", scale_x, scale_y)
        draw_vector_structure(slide, bullets, title_text, scale_x, scale_y)
    elif has_image and not is_auto_image:
        render_split_slide(slide, bullets, image_name, scale_x, scale_y)
    elif is_stats_slide(bullets):
        render_stats_layout(slide, bullets, scale_x, scale_y)
    elif len(bullets) == 3:
        render_three_columns_layout(slide, bullets, scale_x, scale_y)
    elif len(bullets) == 4:
        build_slide_2_cards(slide, bullets, scale_x, scale_y)
    elif has_image:
        render_split_slide(slide, bullets, image_name, scale_x, scale_y)
    else:
        render_default_bullet_layout(slide, bullets, scale_x, scale_y)

def build_presentation(json_data_or_path, output_filename="ranepa_presentation.pptx"):
    """Создает профессиональную презентацию на основе переданных данных слайдов.
    
    Параметры:
    - json_data_or_path: путь к файлу JSON или готовый словарь/список слайдов.
    - output_filename: имя файла для сохранения.
    """
    # Загружаем данные, если передан путь к файлу
    if isinstance(json_data_or_path, str):
        with open(json_data_or_path, "r", encoding="utf-8") as f:
            data = json.load(f)
    else:
        data = json_data_or_path
        
    # Если данные пришли в виде объекта Pydantic, преобразуем в словарь
    if hasattr(data, "dict"):
        data = data.dict()
    if isinstance(data, dict) and "slides" in data:
        slides_list = data["slides"]
    else:
        slides_list = data

    # 1. Инициализируем презентацию (используем template.pptx, если он есть)
    template_path = "template.pptx"
    use_template = os.path.exists(template_path)
    
    if use_template:
        print(f"Используем фирменный шаблон презентации: {template_path}")
        prs = Presentation(template_path)
    else:
        print("Фирменный шаблон template.pptx не найден, создаем стандартную презентацию.")
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    
    for i, slide_data in enumerate(slides_list):
        slide_num = slide_data.get("slide_number", i + 1)
        title_text = slide_data.get("title", "")
        bullets = slide_data.get("bullets", [])
        
        # Если в данных слайда есть структурированные данные графика, генерируем его картинку автоматически
        chart_data_dict = slide_data.get("chart_data", None)
        image_keywords = slide_data.get("image_keywords", None)
        
        if chart_data_dict:
            generate_chart_image(chart_data_dict)
            image_name = chart_data_dict.get("filename")
        elif image_keywords and not slide_data.get("image"):
            # Запускаем динамическую загрузку картинки по ключевым словам, предоставленным агентом
            # Извлекаем префикс темы из имени файла презентации для избежания конфликтов между темами
            base_out = os.path.basename(output_filename)
            topic_prefix = "theme"
            if base_out.startswith("Презентация_"):
                topic_prefix = base_out.replace("Презентация_", "").replace(".pptx", "")
            elif base_out.startswith("test_"):
                topic_prefix = "test"
            image_filename = f"slide_{slide_num}_{topic_prefix}_theme.jpg"
            image_name = download_thematic_image(image_keywords, image_filename)
        else:
            image_name = slide_data.get("image", None)
            
        # Получаем или добавляем слайд
        if use_template and i < len(prs.slides):
            slide = prs.slides[i]
        else:
            blank_layout_idx = 6 if len(prs.slide_layouts) > 6 else len(prs.slide_layouts) - 1
            blank_layout = prs.slide_layouts[blank_layout_idx]
            slide = prs.slides.add_slide(blank_layout)
            
        scale_x = prs.slide_width.inches / 16.0
        scale_y = prs.slide_height.inches / 9.0
        
        # Отрисовка и классификация слайда
        classify_and_render_slide(slide, title_text, bullets, image_name, i, scale_x, scale_y, use_template)

    # 2. Сохраняем готовую презентацию
    os.makedirs("output", exist_ok=True)
    save_path = os.path.join("output", output_filename)
    prs.save(save_path)
    print(f"Презентация сохранена в: {save_path}")
    
    # 3. Копируем файл в системную папку «Документы» пользователя
    try:
        user_profile = os.environ.get("USERPROFILE")
        if user_profile:
            documents_dir = os.path.join(user_profile, "Documents")
            if os.path.exists(documents_dir):
                abs_save_path = os.path.abspath(save_path)
                abs_target_path = os.path.abspath(os.path.join(documents_dir, output_filename))
                shutil.copy2(abs_save_path, abs_target_path)
                print(f"Копия файла успешно сохранена в ваши Документы: {abs_target_path}")
                return abs_target_path
    except Exception as e:
        print(f"Предупреждение: не удалось скопировать в Документы ({e})")
        
    return save_path

if __name__ == "__main__":
    # Тестовый запуск генератора (заглушка данных)
    test_slides = [
        {
            "slide_number": 1,
            "title": "Управление городским бюджетом",
            "bullets": ["Докладчик: Иванов И.И.", "РАНХиГС, Выпуск 2026", "Научный руководитель: Петров П.П."],
            "image": None
        },
        {
            "slide_number": 2,
            "title": "Актуальность проекта",
            "bullets": ["Дефицит бюджета города составляет 15%.", "Необходима оптимизация расходов на ЖКХ.", "Проект предлагает внедрение энергоэффективных систем."],
            "image": None
        }
    ]
    build_presentation(test_slides, "test_presentation.pptx")
