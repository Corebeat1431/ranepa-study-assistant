import os
import sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Reconfigure stdout for Windows console (just in case)
if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8')

def inspect_presentation(path, out_file):
    out_file.write(f"\n=======================================================\n")
    out_file.write(f"ИССЛЕДОВАНИЕ ФАЙЛА: {os.path.basename(path)}\n")
    out_file.write(f"=======================================================\n")
    if not os.path.exists(path):
        out_file.write("Файл не найден!\n")
        return
    
    prs = Presentation(path)
    out_file.write(f"Количество слайдов: {len(prs.slides)}\n")
    out_file.write(f"Ширина слайда: {prs.slide_width.inches:.2f}\"\n")
    out_file.write(f"Высота слайда: {prs.slide_height.inches:.2f}\"\n")
    
    for i, slide in enumerate(prs.slides):
        out_file.write(f"\n--- Слайд {i+1} ---\n")
        
        # Инспектируем макет
        layout = slide.slide_layout
        out_file.write(f"Макет: {layout.name} (Индекс: {prs.slide_layouts.index(layout)})\n")
        
        # Инспектируем фигуры на слайде
        shapes = slide.shapes
        out_file.write(f"Количество объектов на слайде: {len(shapes)}\n")
        
        for shape_idx, shape in enumerate(shapes):
            shape_type = shape.shape_type
            out_file.write(f"  [{shape_idx+1}] Имя: '{shape.name}', Тип: {shape_type}\n")
            
            # Размеры и положение
            left = shape.left.inches if shape.left else 0
            top = shape.top.inches if shape.top else 0
            width = shape.width.inches if shape.width else 0
            height = shape.height.inches if shape.height else 0
            out_file.write(f"      Положение: X={left:.2f}\", Y={top:.2f}\", Размер: W={width:.2f}\", H={height:.2f}\"\n")
            
            # Текст (если есть)
            if shape.has_text_frame:
                tf = shape.text_frame
                text = tf.text.strip()
                if text:
                    out_file.write(f"      Текст (первые 300 символов): \"{text[:300]}\"\n")
                    # Проверяем шрифты в первом параграфе
                    if tf.paragraphs:
                        p = tf.paragraphs[0]
                        if p.runs:
                            run = p.runs[0]
                            font_name = run.font.name
                            font_size = run.font.size.pt if run.font.size else "По умолчанию"
                            font_bold = run.font.bold
                            font_color = run.font.color.rgb if run.font.color and hasattr(run.font.color, 'rgb') else "По умолчанию"
                            out_file.write(f"      Шрифт: {font_name}, Размер: {font_size}, Жирный: {font_bold}, Цвет RGB: {font_color}\n")
            
            # Если это картинка
            if shape_type == MSO_SHAPE_TYPE.PICTURE:
                img_name = shape.image.filename if hasattr(shape.image, 'filename') else 'Встроенный байт-код'
                out_file.write(f"      Это картинка. Файл изображения: {img_name}\n")

if __name__ == "__main__":
    with open("pptx_inspection_results.txt", "w", encoding="utf-8") as out:
        inspect_presentation("Nauchno-populyarnyj-turizm-v-g-Sevastopole.pptx", out)
        inspect_presentation("sablon_prezentacii.pptx", out)
    print("Отчет успешно записан в pptx_inspection_results.txt")
